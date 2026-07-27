"""Vault indexer — text extraction, chunking, FTS5 + FAISS indexing.

Indexes vault files into the shared FTS5 and FAISS indexes used by the
retrieval pipeline. Vault chunk IDs use the prefix "vault:" to distinguish
them from memory chunks.

Chunk ID format: vault:{vault_rel_path}:chunk_{N}
"""

from __future__ import annotations

import csv
import json
import logging
from pathlib import Path

from tesseract.context.circuit_breaker import CircuitBreaker
from tesseract.memory.embeddings import EmbeddingIndex
from tesseract.memory.fts_index import FTSIndex

logger = logging.getLogger(__name__)

try:
    import pdfplumber  # type: ignore[import-untyped]
    _HAS_PDF = True
except ImportError:
    _HAS_PDF = False

try:
    import docx  # type: ignore[import-untyped]
    _HAS_DOCX = True
except ImportError:
    _HAS_DOCX = False

try:
    from pptx import Presentation as _PptxPresentation  # type: ignore[import-untyped]
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False

_CHUNK_SIZE = 400
_CHUNK_OVERLAP = 80
_CSV_MAX_ROWS = 51
_JSON_MAX_CHARS = 5000


class VaultIndexer:
    def __init__(
        self,
        embeddings: EmbeddingIndex,
        fts_index: FTSIndex | None = None,
        log_dir: Path | None = None,
    ) -> None:
        self._embeddings = embeddings
        self._fts_index = fts_index
        self._breaker = CircuitBreaker(
            name="vault_indexer",
            max_failures=3,
            log_dir=log_dir,
        )

    @staticmethod
    def extract_text(vault_path: Path) -> str | None:
        """Extract searchable text from a vault file. Returns None if unsupported."""
        ext = vault_path.suffix.lower()

        if ext in (".md", ".txt"):
            try:
                return vault_path.read_text(encoding="utf-8", errors="replace")
            except OSError:
                logger.warning("Failed to read %s", vault_path)
                return None

        if ext == ".pdf":
            if not _HAS_PDF:
                logger.warning("pdfplumber not installed, cannot index %s", vault_path.name)
                return None
            try:
                with pdfplumber.open(str(vault_path)) as pdf:
                    return "\n".join(page.extract_text() or "" for page in pdf.pages)
            except Exception:
                logger.warning("PDF extraction failed for %s", vault_path.name)
                return None

        if ext in (".csv", ".tsv"):
            try:
                delimiter = "\t" if ext == ".tsv" else ","
                with open(vault_path, encoding="utf-8", errors="replace") as f:
                    reader = csv.reader(f, delimiter=delimiter)
                    rows = [", ".join(row) for _, row in zip(range(_CSV_MAX_ROWS), reader)]
                    return "\n".join(rows)
            except Exception:
                logger.warning("CSV extraction failed for %s", vault_path.name)
                return None

        if ext == ".json":
            try:
                with open(vault_path, encoding="utf-8") as f:
                    data = json.load(f)
                    return json.dumps(data, indent=2)[:_JSON_MAX_CHARS]
            except Exception:
                logger.warning("JSON extraction failed for %s", vault_path.name)
                return None

        if ext == ".docx":
            if not _HAS_DOCX:
                logger.warning("python-docx not installed, cannot index %s", vault_path.name)
                return None
            try:
                doc = docx.Document(str(vault_path))
                return "\n".join(p.text for p in doc.paragraphs)
            except Exception:
                logger.warning("DOCX extraction failed for %s", vault_path.name)
                return None

        if ext == ".pptx":
            if not _HAS_PPTX:
                logger.warning("python-pptx not installed, cannot index %s", vault_path.name)
                return None
            try:
                prs = _PptxPresentation(str(vault_path))
                parts: list[str] = []
                for idx, slide in enumerate(prs.slides, start=1):
                    slide_lines: list[str] = [f"# Slide {idx}"]
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                text = "".join(run.text for run in para.runs).strip()
                                if text:
                                    slide_lines.append(text)
                    notes = slide.notes_slide if slide.has_notes_slide else None
                    if notes is not None and notes.notes_text_frame is not None:
                        notes_text = notes.notes_text_frame.text.strip()
                        if notes_text:
                            slide_lines.append("Notes:")
                            slide_lines.append(notes_text)
                    parts.append("\n".join(slide_lines))
                return "\n\n".join(parts)
            except Exception:
                logger.warning("PPTX extraction failed for %s", vault_path.name)
                return None

        # Attempt raw UTF-8 read for unknown types
        try:
            return vault_path.read_text(encoding="utf-8", errors="strict")
        except (UnicodeDecodeError, OSError):
            return None

    @staticmethod
    def chunk_content(
        content: str,
        chunk_size: int = _CHUNK_SIZE,
        overlap: int = _CHUNK_OVERLAP,
    ) -> list[str]:
        """Split content into overlapping chunks by word count."""
        words = content.split()
        if not words:
            return []

        chunks: list[str] = []
        start = 0
        while start < len(words):
            end = start + chunk_size
            chunk = " ".join(words[start:end])
            if chunk.strip():
                chunks.append(chunk)
            start = end - overlap

        return chunks

    @staticmethod
    def chunk_id(vault_rel_path: str, index: int) -> str:
        """Build a chunk ID from vault-relative path and chunk index."""
        return f"vault:{vault_rel_path}:chunk_{index}"

    async def index_vault_file(
        self,
        vault_rel_path: str,
        title: str,
        vault_abs_path: Path,
    ) -> int:
        """Extract text, chunk, and index into FTS5 + FAISS. Returns chunk count."""
        await self.remove_vault_file(vault_rel_path)

        content = self.extract_text(vault_abs_path)
        if not content:
            logger.info("No extractable text from %s", vault_rel_path)
            return 0

        chunks = self.chunk_content(content)
        if not chunks:
            return 0

        indexed = 0
        for i, chunk_text in enumerate(chunks):
            cid = self.chunk_id(vault_rel_path, i)

            if self._fts_index is not None:
                try:
                    self._fts_index.add(cid, title, chunk_text)
                except Exception:
                    logger.warning("FTS add failed for %s", cid)

            if not self._breaker.is_tripped:
                try:
                    await self._embeddings.add(cid, chunk_text)
                except Exception as e:
                    self._breaker.record_failure(str(e))
                    logger.warning("Embedding failed for %s: %s", cid, e)

            indexed += 1

        logger.info("Indexed %d chunks from %s", indexed, vault_rel_path)
        return indexed

    async def remove_vault_file(self, vault_rel_path: str) -> int:
        """Remove all chunks for a vault file from FTS and FAISS."""
        prefix = f"vault:{vault_rel_path}:chunk_"
        removed = 0

        if self._fts_index is not None:
            all_ids = self._fts_index.all_ids()
            for fts_id in all_ids:
                if fts_id.startswith(prefix):
                    self._fts_index.delete(fts_id)
                    removed += 1

        # Remove from FAISS — scan id_map for matching keys. Use the
        # public `snapshot_ids()` so the lookup acquires the embeddings
        # lock instead of reaching into private `_id_to_pos` (audit-fix M5).
        try:
            for mem_id in self._embeddings.snapshot_ids():
                if mem_id.startswith(prefix):
                    self._embeddings.remove(mem_id)
        except Exception:
            logger.warning("FAISS removal failed for %s", vault_rel_path)

        return removed
