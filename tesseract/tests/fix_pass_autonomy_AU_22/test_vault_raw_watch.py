"""AU-22 VaultRawWatchJob — 14 test cases.

Coverage:

1.  suggest_raw_filing_path uses YYYYMMDD/ granularity
2.  empty-tick short-circuit (no folders, no event)
3.  non-conforming folder names skipped + reported
4.  mixed-type batch — pdf + docx + txt auto, .url + oversized to ASK
5.  cursor jsonl dedup by SHA (rename doesn't re-trigger)
6.  prior-failure routes to ASK, never auto-retries
7.  mode: ask_all forces every file through ASK
8.  Auto-allowlist promotes .url to AUTO when prefix matches
9.  .pptx extraction (slides + notes joined)
10. .url shortcut extractor reads the URL field
11. apply_ask_batch — operator approval triggers ingest + cursor row
12. apply_ask_batch — denial recorded without ingest
13. Workspace event has the right kind + payload shape
14. Files never move out of their YYYYMMDD/ folder
"""

from __future__ import annotations

import asyncio
import json
import zipfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import pytest

from tesseract.memory._shortcut_extractor import extract_url
from tesseract.memory.vault_indexer import VaultIndexer, _HAS_PPTX
from tesseract.memory.vault_manager import VaultManager
from tesseract.scheduler.tasks.vault_raw_watch import (
    DEFAULT_MAX_AUTO_SIZE_MB,
    VaultRawWatchJob,
    apply_ask_batch,
)
from tesseract.scheduler.types import JobContext
from tesseract.workspace_events.events import EventStore


@pytest.fixture
def isolated_home(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> Path:
    monkeypatch.setenv("TESSERACT_HOME", str(tmp_path))
    return tmp_path


def _seed_vault(home: Path) -> VaultManager:
    vm = VaultManager(vault_root=home / "vault")
    (home / "vault" / "raw").mkdir(parents=True, exist_ok=True)
    return vm


def _make_ctx(
    *,
    home: Path,
    vm: VaultManager,
    event_store: EventStore,
    raw_watch: dict[str, Any] | None = None,
    fired_at: datetime | None = None,
) -> JobContext:
    cfg: dict[str, Any] = {
        "vault_manager": vm,
        "vault_indexer": None,
        "cursor_path": str(home / "autonomy" / "vault-raw-cursors.jsonl"),
        "event_store": event_store,
    }
    if raw_watch is not None:
        cfg["raw_watch"] = raw_watch
    return JobContext(
        job_name="vault_raw_watch",
        config=cfg,
        fired_at=fired_at or datetime(2026, 5, 18, 4, 0, 0, tzinfo=timezone.utc),
    )


def _read_cursor(home: Path) -> list[dict[str, Any]]:
    path = home / "autonomy" / "vault-raw-cursors.jsonl"
    if not path.exists():
        return []
    return [json.loads(line) for line in path.read_text(encoding="utf-8").splitlines() if line.strip()]


# 1 ─────────────────────────────────────────────────────────────────


def test_suggest_raw_filing_path_uses_yyyymmdd(isolated_home: Path) -> None:
    vm = VaultManager(vault_root=isolated_home / "vault")
    suggestion = vm.suggest_raw_filing_path("Some Research.PDF")
    # Today's date in UTC, eight digits
    today = datetime.now(timezone.utc).strftime("%Y%m%d")
    assert suggestion == f"raw/{today}/some-research.pdf"


# 2 ─────────────────────────────────────────────────────────────────


def test_empty_tick_short_circuits(isolated_home: Path) -> None:
    vm = _seed_vault(isolated_home)
    store = EventStore(isolated_home / "logs")
    ctx = _make_ctx(home=isolated_home, vm=vm, event_store=store)
    result = asyncio.run(VaultRawWatchJob().run(ctx))
    assert result.ok
    assert result.payload["folders_scanned"] == 0
    assert result.payload["auto_ingested"] == []
    assert store.events_path.exists() is False  # no event emitted


# 3 ─────────────────────────────────────────────────────────────────


def test_nonconforming_folder_names_skipped(isolated_home: Path) -> None:
    vm = _seed_vault(isolated_home)
    (isolated_home / "vault" / "raw" / "2026-04").mkdir()
    (isolated_home / "vault" / "raw" / "research").mkdir()
    (isolated_home / "vault" / "raw" / "20260518").mkdir()
    (isolated_home / "vault" / "raw" / "20260518" / "note.txt").write_text("hello", encoding="utf-8")
    store = EventStore(isolated_home / "logs")
    ctx = _make_ctx(home=isolated_home, vm=vm, event_store=store)
    result = asyncio.run(VaultRawWatchJob().run(ctx))
    assert result.ok
    assert result.payload["folders_scanned"] == 1
    assert sorted(result.payload["skipped_nonconforming"]) == ["2026-04", "research"]


# 4 ─────────────────────────────────────────────────────────────────


def test_mixed_batch_routes_correctly(isolated_home: Path) -> None:
    vm = _seed_vault(isolated_home)
    folder = isolated_home / "vault" / "raw" / "20260518"
    folder.mkdir()
    (folder / "small.txt").write_text("tiny note", encoding="utf-8")
    (folder / "report.md").write_text("# Report\nLine", encoding="utf-8")
    huge = folder / "huge.bin"
    huge.write_bytes(b"x" * (DEFAULT_MAX_AUTO_SIZE_MB * 1024 * 1024 + 10))
    (folder / "link.url").write_text(
        "[InternetShortcut]\nURL=https://example.com/article\n", encoding="utf-8"
    )

    store = EventStore(isolated_home / "logs")
    ctx = _make_ctx(home=isolated_home, vm=vm, event_store=store)
    result = asyncio.run(VaultRawWatchJob().run(ctx))
    assert result.ok
    # Auto = small.txt + report.md; ASK = huge.bin + link.url
    auto = set(result.payload["auto_ingested"])
    assert auto == {"20260518/small.txt", "20260518/report.md"}
    assert result.payload["ask_queued_count"] == 2

    cursor_rows = _read_cursor(isolated_home)
    auto_rows = [r for r in cursor_rows if r["decision"] == "auto"]
    assert len(auto_rows) == 2
    for row in auto_rows:
        assert row["ingest_status"] == "ingested"

    # Files MUST still exist in their date folder
    assert (folder / "small.txt").exists()
    assert (folder / "report.md").exists()
    assert (folder / "huge.bin").exists()
    assert (folder / "link.url").exists()


# 5 ─────────────────────────────────────────────────────────────────


def test_cursor_dedups_by_sha_not_path(isolated_home: Path) -> None:
    vm = _seed_vault(isolated_home)
    folder = isolated_home / "vault" / "raw" / "20260518"
    folder.mkdir()
    (folder / "doc.txt").write_text("identical bytes", encoding="utf-8")
    store = EventStore(isolated_home / "logs")
    ctx = _make_ctx(home=isolated_home, vm=vm, event_store=store)
    result1 = asyncio.run(VaultRawWatchJob().run(ctx))
    assert "20260518/doc.txt" in result1.payload["auto_ingested"]

    # Rename the file — same SHA, different path. Watcher should NOT re-ingest.
    (folder / "doc.txt").rename(folder / "renamed.txt")
    result2 = asyncio.run(VaultRawWatchJob().run(ctx))
    assert result2.payload["auto_ingested"] == []
    assert result2.payload["ask_queued_count"] == 0


# 6 ─────────────────────────────────────────────────────────────────


def test_prior_failure_never_retries_until_sha_changes(isolated_home: Path) -> None:
    """AU-22 contract: a `(folder, sha)` row with `ingest_status: failed`
    blocks ALL re-proposal (auto AND ask) until the operator replaces
    the file (producing a new SHA). The dedup at `_read_cursor_keys`
    handles this — any prior row marks the pair as seen."""
    vm = _seed_vault(isolated_home)
    folder = isolated_home / "vault" / "raw" / "20260518"
    folder.mkdir()
    target = folder / "report.md"
    target.write_text("# Report", encoding="utf-8")
    import hashlib

    sha_old = hashlib.sha256(target.read_bytes()).hexdigest()
    cursor = isolated_home / "autonomy" / "vault-raw-cursors.jsonl"
    cursor.parent.mkdir(parents=True, exist_ok=True)
    cursor.write_text(
        json.dumps(
            {
                "folder": "20260518",
                "relpath": "20260518/report.md",
                "sha256": sha_old,
                "ingest_status": "failed",
                "decision": "auto",
                "reason": "simulated extractor crash",
            }
        )
        + "\n",
        encoding="utf-8",
    )

    store = EventStore(isolated_home / "logs")
    ctx = _make_ctx(home=isolated_home, vm=vm, event_store=store)

    # Tick 1: same content (same SHA) → blocked by dedup
    result = asyncio.run(VaultRawWatchJob().run(ctx))
    assert result.payload["auto_ingested"] == []
    assert result.payload["ask_queued_count"] == 0

    # Operator "replaces" the file with new content → new SHA → re-proposed
    target.write_text("# Report (revised)\nNow with content", encoding="utf-8")
    sha_new = hashlib.sha256(target.read_bytes()).hexdigest()
    assert sha_new != sha_old
    result2 = asyncio.run(VaultRawWatchJob().run(ctx))
    assert "20260518/report.md" in result2.payload["auto_ingested"]


# 7 ─────────────────────────────────────────────────────────────────


def test_mode_ask_all_forces_every_file(isolated_home: Path) -> None:
    vm = _seed_vault(isolated_home)
    folder = isolated_home / "vault" / "raw" / "20260518"
    folder.mkdir()
    (folder / "a.txt").write_text("a", encoding="utf-8")
    (folder / "b.txt").write_text("b", encoding="utf-8")
    store = EventStore(isolated_home / "logs")
    ctx = _make_ctx(home=isolated_home, vm=vm, event_store=store, raw_watch={"mode": "ask_all"})
    result = asyncio.run(VaultRawWatchJob().run(ctx))
    assert result.payload["auto_ingested"] == []
    assert result.payload["ask_queued_count"] == 2


# 8 ─────────────────────────────────────────────────────────────────


def test_auto_url_allowlist_promotes_shortcut_to_auto(isolated_home: Path) -> None:
    vm = _seed_vault(isolated_home)
    folder = isolated_home / "vault" / "raw" / "20260518"
    folder.mkdir()
    (folder / "trusted.url").write_text(
        "[InternetShortcut]\nURL=https://docs.example.com/page\n", encoding="utf-8"
    )
    (folder / "other.url").write_text(
        "[InternetShortcut]\nURL=https://random.example.com/article\n", encoding="utf-8"
    )
    store = EventStore(isolated_home / "logs")
    ctx = _make_ctx(
        home=isolated_home,
        vm=vm,
        event_store=store,
        raw_watch={"auto_url_allowlist": ["https://docs.example.com/"]},
    )
    result = asyncio.run(VaultRawWatchJob().run(ctx))
    assert "20260518/trusted.url" in result.payload["auto_ingested"]
    assert result.payload["ask_queued_count"] == 1


# 9 ─────────────────────────────────────────────────────────────────


@pytest.mark.skipif(not _HAS_PPTX, reason="python-pptx not installed")
def test_pptx_extractor_returns_slides_and_notes(isolated_home: Path) -> None:
    """Use python-pptx to build a real .pptx so we exercise the actual
    extractor branch end-to-end."""
    from pptx import Presentation

    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[5])
    title = slide.shapes.title
    if title is not None:
        title.text = "Hello slide"
    slide.notes_slide.notes_text_frame.text = "Speaker notes here"
    path = isolated_home / "deck.pptx"
    pres.save(str(path))

    text = VaultIndexer.extract_text(path) or ""
    assert "Hello slide" in text
    assert "Speaker notes here" in text


def test_pptx_branch_handles_missing_dep(monkeypatch: pytest.MonkeyPatch, tmp_path: Path) -> None:
    """When python-pptx is missing the extractor must return None, not raise."""
    import tesseract.memory.vault_indexer as mod

    fake = tmp_path / "deck.pptx"
    # Build a minimal valid-looking zip so vault_indexer's earlier guards
    # don't trip — the .pptx branch only fires after extension match.
    with zipfile.ZipFile(fake, "w") as z:
        z.writestr("[Content_Types].xml", "<x/>")
    monkeypatch.setattr(mod, "_HAS_PPTX", False)
    assert VaultIndexer.extract_text(fake) is None


# 10 ────────────────────────────────────────────────────────────────


def test_shortcut_extractor_reads_url(isolated_home: Path) -> None:
    path = isolated_home / "test.url"
    path.write_text("[InternetShortcut]\nURL=https://example.org/path\n", encoding="utf-8")
    target = extract_url(path)
    assert target.url == "https://example.org/path"


def test_shortcut_extractor_lnk_byte_scan(isolated_home: Path) -> None:
    path = isolated_home / "test.lnk"
    # Real .lnk files are binary; embed a URL in a blob so the fallback scan
    # recovers it.
    path.write_bytes(b"\x00\x01\x02 some prelude https://example.org/lnk-target trailing\x00")
    target = extract_url(path)
    assert target.url == "https://example.org/lnk-target"


# 11, 12 ────────────────────────────────────────────────────────────


def test_apply_ask_batch_ingests_on_approve(isolated_home: Path) -> None:
    vm = _seed_vault(isolated_home)
    folder = isolated_home / "vault" / "raw" / "20260518"
    folder.mkdir()
    target = folder / "doc.md"
    target.write_text("# title\nbody", encoding="utf-8")
    import hashlib

    sha = hashlib.sha256(target.read_bytes()).hexdigest()

    files = [
        {
            "folder": "20260518",
            "relpath": "20260518/doc.md",
            "sha256": sha,
            "size_bytes": target.stat().st_size,
        }
    ]
    cursor = isolated_home / "autonomy" / "vault-raw-cursors.jsonl"
    summary = asyncio.run(
        apply_ask_batch(
            files=files,
            decisions={},
            vault_manager=vm,
            indexer=None,
            cursor_path=cursor,
        )
    )
    assert summary["ingested"] == ["20260518/doc.md"]
    rows = _read_cursor(isolated_home)
    assert any(r["decision"] == "ask" and r["ingest_status"] == "ingested" for r in rows)
    # Sidecar should land alongside the file
    sidecar = folder / "doc.md.meta.yaml"
    assert sidecar.exists()


def test_apply_ask_batch_records_deny(isolated_home: Path) -> None:
    vm = _seed_vault(isolated_home)
    folder = isolated_home / "vault" / "raw" / "20260518"
    folder.mkdir()
    target = folder / "doc.md"
    target.write_text("body", encoding="utf-8")
    import hashlib

    sha = hashlib.sha256(target.read_bytes()).hexdigest()
    files = [
        {
            "folder": "20260518",
            "relpath": "20260518/doc.md",
            "sha256": sha,
            "size_bytes": target.stat().st_size,
        }
    ]
    cursor = isolated_home / "autonomy" / "vault-raw-cursors.jsonl"
    summary = asyncio.run(
        apply_ask_batch(
            files=files,
            decisions={"20260518/doc.md": "denied"},
            vault_manager=vm,
            indexer=None,
            cursor_path=cursor,
        )
    )
    assert summary["denied"] == ["20260518/doc.md"]
    assert summary["ingested"] == []
    # No sidecar should be written on deny
    assert not (folder / "doc.md.meta.yaml").exists()
    rows = _read_cursor(isolated_home)
    assert any(r["decision"] == "ask" and r["ingest_status"] == "denied" for r in rows)


# 13 ────────────────────────────────────────────────────────────────


def test_ask_event_payload_shape(isolated_home: Path) -> None:
    vm = _seed_vault(isolated_home)
    folder = isolated_home / "vault" / "raw" / "20260518"
    folder.mkdir()
    huge = folder / "huge.bin"
    huge.write_bytes(b"x" * (DEFAULT_MAX_AUTO_SIZE_MB * 1024 * 1024 + 1))
    store = EventStore(isolated_home / "logs")
    ctx = _make_ctx(home=isolated_home, vm=vm, event_store=store)
    asyncio.run(VaultRawWatchJob().run(ctx))
    events = store.list_events(kinds=("vault_raw_ingest_batch",))
    assert len(events) == 1
    ev = events[0]
    assert ev.kind == "vault_raw_ingest_batch"
    files = ev.payload["files"]
    assert len(files) == 1
    f = files[0]
    assert f["relpath"] == "20260518/huge.bin"
    assert "ask_reason" in f
    assert f["suggested_path"] == "raw/20260518/huge.bin"
    assert ev.payload["folders"] == ["20260518"]


# 14 ────────────────────────────────────────────────────────────────


def test_files_never_move_out_of_date_folder(isolated_home: Path) -> None:
    vm = _seed_vault(isolated_home)
    folder = isolated_home / "vault" / "raw" / "20260518"
    folder.mkdir()
    (folder / "note.md").write_text("# note", encoding="utf-8")
    store = EventStore(isolated_home / "logs")
    ctx = _make_ctx(home=isolated_home, vm=vm, event_store=store)
    asyncio.run(VaultRawWatchJob().run(ctx))
    # File must still be in the same folder
    assert (folder / "note.md").exists()
    # And meta sidecar landed alongside
    assert (folder / "note.md.meta.yaml").exists()
    # No copy under another vault category
    assert not any(p.is_dir() and p.name != "raw" for p in (isolated_home / "vault").iterdir())


# Disabled-mode short-circuit -----------------------------------------


def test_disabled_flag_short_circuits(isolated_home: Path) -> None:
    vm = _seed_vault(isolated_home)
    folder = isolated_home / "vault" / "raw" / "20260518"
    folder.mkdir()
    (folder / "note.md").write_text("# n", encoding="utf-8")
    store = EventStore(isolated_home / "logs")
    ctx = _make_ctx(home=isolated_home, vm=vm, event_store=store, raw_watch={"enabled": False})
    result = asyncio.run(VaultRawWatchJob().run(ctx))
    assert result.ok
    assert result.detail == "raw_watch disabled"
    assert _read_cursor(isolated_home) == []
